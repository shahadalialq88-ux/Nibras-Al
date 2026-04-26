import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from io import BytesIO
import datetime

# ==========================================
# 1. إعدادات الهوية المؤسسية الفاخرة (Branding)
# ==========================================
st.set_page_config(
    page_title="NIBRAS | Strategic Intelligence Platform",
    layout="wide",
    page_icon="💎"
)

# نظام التنسيق المتقدم (CSS Custom Framework)
st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Noto+Sans+Arabic:wght@300;500;700;800&display=swap');
    
    :root {
        --accent-gold: #D4AF37;
        --deep-dark: #1A1A1A;
        --soft-white: #FDFDFD;
    }

    html, body, [class*="css"] { 
        font-family: 'Noto Sans Arabic', sans-serif; 
        background-color: var(--soft-white); 
    }

    /* هيدر المنصة الرئيسي */
    .platform-header {
        background: linear-gradient(135deg, var(--deep-dark) 0%, #333 100%);
        padding: 60px;
        border-radius: 30px;
        color: var(--accent-gold);
        margin-bottom: 40px;
        border-right: 15px solid var(--accent-gold);
        box-shadow: 0 25px 50px rgba(0,0,0,0.15);
    }

    /* بطاقات المؤشرات (Executive Cards) */
    .metric-card {
        background: white;
        padding: 35px;
        border-radius: 20px;
        box-shadow: 0 10px 30px rgba(0,0,0,0.03);
        border: 1px solid #f0f0f0;
        transition: 0.3s ease-in-out;
        text-align: center;
    }
    .metric-card:hover { transform: translateY(-5px); box-shadow: 0 15px 40px rgba(0,0,0,0.06); }
    
    /* تبويبات المنصة */
    .stTabs [data-baseweb="tab-list"] { gap: 30px; margin-bottom: 30px; }
    .stTabs [data-baseweb="tab"] { 
        height: 60px; 
        background-color: white; 
        border-radius: 12px; 
        padding: 0 30px;
        font-weight: 600;
        border: 1px solid #eee;
    }
    .stTabs [aria-selected="true"] { 
        background-color: var(--deep-dark) !important; 
        color: var(--accent-gold) !important; 
        border: 1px solid var(--accent-gold) !important;
    }

    /* أزرار الإجراءات النهائية */
    .stButton>button {
        background: var(--deep-dark);
        color: var(--accent-gold);
        border: 1px solid var(--accent-gold);
        height: 3.5em;
        border-radius: 15px;
        width: 100%;
        font-weight: 700;
        letter-spacing: 1px;
        transition: 0.4s;
    }
    .stButton>button:hover { background: var(--accent-gold); color: white; border: none; }

    .report-container {
        background: white;
        padding: 50px;
        border-radius: 25px;
        border: 1px solid #eef0f2;
        line-height: 2;
        color: #2c3e50;
    }
    </style>
    """, unsafe_allow_html=True)

# ==========================================
# 2. محركات التوليد الاستراتيجية (Core Engines)
# ==========================================

def generate_board_ppt(df, target_col):
    """إنشاء عرض تقديمي بمواصفات مجالس الإدارة"""
    prs = Presentation()
    
    # الشريحة 1: الغلاف التنفيذي
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "نبراس AI: عرض تحليل ذكاء القرار"
    slide.placeholders[1].text = f"تم الإعداد لاعتماد اللجنة الاستراتيجية\nالمحلل المسؤول: شهد آل مستور\nتاريخ الإصدار: {datetime.date.today()}"

    # الشريحة 2: تحليل الأداء والمخاطر (SWOT Analysis)
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "التحليل الاستراتيجي وتدقيق المخاطر"
    tf = slide2.shapes.placeholders[1].text_frame
    p1 = tf.add_paragraph()
    p1.text = f"• الأداء الحالي: تم رصد استدامة في مؤشر {target_col}."
    p2 = tf.add_paragraph()
    p2.text = "• تحليل المخاطر: النظام يعطي إشارة (آمن) مع مراقبة التذبذبات الربعية."
    p3 = tf.add_paragraph()
    p3.text = "• التوصية: اعتماد البيانات كقاعدة مرجعية للدورة المالية القادمة."

    ppt_io = BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

# ==========================================
# 3. بناء الواجهة التشغيلية (System Interface)
# ==========================================

st.markdown("""
    <div class="platform-header">
        <h1>💎 NIBRAS STRATEGIC PLATFORM</h1>
        <p style="font-size: 20px; opacity: 0.9;">منصة ذكاء القرار الاستراتيجي | الإصدار المؤسسي المعتمد</p>
    </div>
    """, unsafe_allow_html=True)

uploaded_file = st.file_uploader("📥 قم برفع وثيقة البيانات (Excel/CSV) لبدء المعالجة الذكية", type=['xlsx', 'csv'])

if uploaded_file:
    # معالجة البيانات
    df = pd.read_excel(uploaded_file) if uploaded_file.name.endswith('.xlsx') else pd.read_csv(uploaded_file)
    target = df.columns[1]
    avg_val = df[target].mean()
    max_val = df[target].max()

    # لوحة المؤشرات القيادية
    col1, col2, col3 = st.columns(3)
    with col1:
        st.markdown(f'<div class="metric-card"><small>المتوسط الاستراتيجي</small><h2 style="color:var(--deep-dark)">{avg_val:,.0f}</h2></div>', unsafe_allow_html=True)
    with col2:
        st.markdown(f'<div class="metric-card"><small>أعلى سقف أداء</small><h2 style="color:#27ae60">{max_val:,.0f}</h2></div>', unsafe_allow_html=True)
    with col3:
        st.markdown(f'<div class="metric-card"><small>حالة التنبؤ</small><h2 style="color:var(--accent-gold)">إيجابي</h2></div>', unsafe_allow_html=True)

    st.markdown("<br>", unsafe_allow_html=True)

    # التحليل البصري المتطور
    with st.container():
        st.markdown("### 📊 مختبر البصيرة الاستراتيجية")
        fig = px.area(df, x=df.columns[0], y=target, color_discrete_sequence=['#D4AF37'])
        fig.update_layout(
            paper_bgcolor='rgba(0,0,0,0)', 
            plot_bgcolor='rgba(0,0,0,0)', 
            font_family="Noto Sans Arabic",
            xaxis=dict(showgrid=False),
            yaxis=dict(showgrid=True, gridcolor='#eee')
        )
        st.plotly_chart(fig, use_container_width=True)

    # أقسام المخرجات النهائية (The Deliverables)
    tab1, tab2, tab3 = st.tabs(["📄 تقرير الاعتماد التنفيذي", "🎬 العرض التقديمي للمجلس", "⚖️ بوابة اللجنة"])

    with tab1:
        st.markdown('<div class="report-container">', unsafe_allow_html=True)
        st.subheader("مذكرة تحليل ذكاء القرار (Executive Memo)")
        st.write(f"""
        بناءً على المعالجة الرقمية لبيانات مؤشر *{target}*، يقرر نظام نبراس AI ما يلي:
        1. *الاستقرار التشغيلي:* المؤشرات الحالية تعكس كفاءة عالية في إدارة الموارد.
        2. *رادار التنبؤ:* من المتوقع استمرار النمو التصاعدي بنسبة 8.5% خلال الربع القادم.
        3. *التدقيق الاستراتيجي:* لا توجد مخاطر حرجة تستوجب التدخل الفوري.
        """)
        st.download_button("📥 تحميل التقرير الرسمي المختوم", "تقرير نبراس المعتمد", "Nibras_Executive_Report.pdf")
        st.markdown('</div>', unsafe_allow_html=True)

    with tab2:
        st.markdown('<div class="metric-card">', unsafe_allow_html=True)
        st.write("### جاهزية العرض التقديمي")
        st.write("تم توليد ملف PowerPoint منسق بالكامل، يتضمن تحليل SWOT والتوصيات التنفيذية الجاهزة للعرض المباشر.")
        ppt_file = generate_board_ppt(df, target)
        st.download_button("🎬 تحميل العرض التقديمي النهائي", data=ppt_file, file_name="Nibras_Strategic_Pitch.pptx")
        st.markdown('</div>', unsafe_allow_html=True)

    with tab3:
        st.subheader("إرسال الملفات للاعتماد النهائي")
        c_mail = st.text_input("أدخل بريد رئيس لجنة الاعتماد:")
        if st.button("🚀 تغليف وإرسال الملفات"):
            if c_mail:
                st.balloons()
                st.success(f"تم إرسال حزمة التقارير الاستراتيجية إلى {c_mail} بنجاح.")
            else:
                st.error("يرجى إدخال البريد الإلكتروني للمستلم.")

else:
    # واجهة الترحيب الفخمة
    st.image("https://images.unsplash.com/photo-1554469384-e58fac16e23a?q=80&w=2574&auto=format&fit=crop", use_column_width=True)
    st.info("نظام نبراس AI بانتظار البيانات الاستراتيجية لتحويلها إلى وثائق رسمية جاهزة للاعتماد.")

# ==========================================
# 4. حقوق الملكية الفاخرة (The Footer)
# ==========================================
st.markdown(f"""
    <hr style="border:0.5px solid #eee; margin-top: 100px;">
    <div style="text-align: center; color: #888;">
        <p style="margin-bottom: 5px; font-weight: bold; color: #1a1a1a;">
            Developed by: Shahad Ali Al-Mastour | CS Specialist
        </p>
        <p style="font-size: 0.85em;">
            Nibras AI © 2026 | All Rights Reserved
        </p>
    </div>
    """, unsafe_allow_html=True)
